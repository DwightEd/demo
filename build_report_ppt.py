"""Rewrite progress-report PPTX per PPT_重写规划_2026-06-08.md.
Rules: no parenthetical explanations; no v1/v2 codenames (use 提示形式+判分规则);
every method states how it is computed; data goes into tables; trained probe included.
Figures use Chinese labels (Microsoft YaHei).
"""
import os
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei"]
plt.rcParams["axes.unicode_minus"] = False
plt.rcParams.update({"figure.dpi": 150, "font.size": 11, "axes.grid": True,
                     "grid.alpha": 0.3, "axes.axisbelow": True})
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = r"C:\Users\613\Desktop\AAAI2026"
FIG = os.path.join(OUT, "fig"); os.makedirs(FIG, exist_ok=True)
BLUE, RED, GRAY, GREEN, ORANGE = "#2f6fb0", "#c0392b", "#9aa0a6", "#2e8b57", "#e08e0b"


def fig1():
    labels = ["跨题随机", "按题留出整体", "按题留出同题配对", "链长基线"]
    vals = [0.898, 0.715, 0.715, 0.585]; cols = [RED, BLUE, BLUE, GRAY]
    fig, ax = plt.subplots(figsize=(6.2, 3.6)); b = ax.bar(labels, vals, color=cols)
    for r, v in zip(b, vals): ax.text(r.get_x()+r.get_width()/2, v+0.008, f"{v:.3f}", ha="center", fontsize=10)
    ax.axhline(0.5, ls="--", c="k", lw=0.8); ax.set_ylim(0.45, 0.95)
    ax.set_ylabel("AUROC"); ax.set_title("训练探针:四种评测方式的 AUROC")
    ax.text(1.0, 0.84, "高出 0.18", ha="center", color=RED, fontsize=11)
    fig.tight_layout(); p = os.path.join(FIG, "f1.png"); fig.savefig(p); plt.close(); return p


def fig2():
    fr = [0.1, 0.3, 0.5, 0.7, 0.9]
    within = [0.554, 0.576, 0.625, 0.691, 0.687]; cross = [0.833, 0.824, 0.844, 0.864, 0.896]
    fig, ax = plt.subplots(figsize=(6.2, 3.6))
    ax.plot(fr, within, "-o", c=BLUE, label="题内配对"); ax.plot(fr, cross, "-s", c=RED, label="跨题")
    ax.axhline(0.5, ls="--", c="k", lw=0.8); ax.set_ylim(0.45, 0.95)
    ax.set_xlabel("推理链中的相对位置  前段 → 后段"); ax.set_ylabel("AUROC")
    ax.set_title("错误检测力沿推理链后移"); ax.legend(fontsize=10, loc="center left")
    fig.tight_layout(); p = os.path.join(FIG, "f2.png"); fig.savefig(p); plt.close(); return p


def fig3():
    names = ["参与度 PR", "激活熵 AE", "Mahalanobis", "流形外能量 SPE", "训练线性探针", "低秩25维探针", "多信号集成"]
    vals = [0.60, 0.58, 0.65, 0.68, 0.71, 0.73, 0.76]
    fig, ax = plt.subplots(figsize=(6.2, 3.6)); y = np.arange(len(names))
    ax.barh(y, vals, color=BLUE); ax.set_yticks(y); ax.set_yticklabels(names, fontsize=10)
    for i, v in enumerate(vals): ax.text(v+0.003, i, f"{v:.2f}", va="center", fontsize=9)
    ax.axvspan(0.68, 0.76, color=GRAY, alpha=0.25, label="上限带 0.68–0.76")
    ax.axvline(0.5, ls="--", c="k", lw=0.8); ax.set_xlim(0.5, 0.82)
    ax.set_xlabel("题内配对 AUROC"); ax.set_title("多种方法收敛到同一区间"); ax.legend(fontsize=9, loc="lower right")
    fig.tight_layout(); p = os.path.join(FIG, "f3.png"); fig.savefig(p); plt.close(); return p


def fig4():
    groups = ["自定义零样本\n宽松判分", "标准五样本\n宽松判分", "标准五样本\n严格判分"]
    mid = [0.628, 0.657, 0.830]; alll = [0.646, 0.647, 0.815]; deep = [0.645, 0.639, 0.811]
    x = np.arange(len(groups)); w = 0.26
    fig, ax = plt.subplots(figsize=(6.4, 3.7))
    ax.bar(x-w, mid, w, label="中层", color=BLUE); ax.bar(x, alll, w, label="全层", color=GREEN)
    ax.bar(x+w, deep, w, label="深层", color=ORANGE)
    for xi, v in zip(x-w, mid): ax.text(xi, v+0.006, f"{v:.2f}", ha="center", fontsize=9)
    ax.axhline(0.5, ls="--", c="k", lw=0.8); ax.set_ylim(0.45, 0.9)
    ax.set_xticks(x); ax.set_xticklabels(groups, fontsize=9)
    ax.set_ylabel("题内配对 AUROC"); ax.set_title("Mahalanobis 距离在三种设置下的检测力")
    ax.legend(fontsize=9, loc="upper left")
    fig.tight_layout(); p = os.path.join(FIG, "f4.png"); fig.savefig(p); plt.close(); return p


def fig5():
    feats = [("mahal 中层", 0.830), ("mahal 全层", 0.815), ("mahal 深层", 0.811),
             ("激活熵 深层", 0.528), ("输出熵", 0.515), ("链长", 0.486),
             ("激活熵 中层", 0.410), ("参与度 中层", 0.396)]
    names = [f for f, _ in feats]; vals = [v for _, v in feats]
    cols = [RED if v >= 0.8 else GRAY for v in vals]
    fig, ax = plt.subplots(figsize=(6.4, 3.8)); y = np.arange(len(names))[::-1]
    ax.barh(y, vals, color=cols); ax.set_yticks(y); ax.set_yticklabels(names, fontsize=10)
    for yi, v in zip(y, vals):
        ax.text(v+0.005 if v > 0.5 else v-0.005, yi, f"{v:.3f}", va="center",
                ha="left" if v > 0.5 else "right", fontsize=9)
    ax.axvline(0.5, ls="--", c="k", lw=0.9); ax.set_xlim(0.35, 0.88)
    ax.set_xlabel("题内配对 AUROC,大于 0.5 表示错误链更高"); ax.set_title("标准五样本加严格判分下各特征的检测力")
    fig.tight_layout(); p = os.path.join(FIG, "f5.png"); fig.savefig(p); plt.close(); return p


figs = [fig1(), fig2(), fig3(), fig4(), fig5()]
print("figures done")

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; FONT = "微软雅黑"


def add_text(slide, l, t, w, h, lines):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        txt, sz, bd, cl = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = txt; f = r.font; f.size = Pt(sz); f.bold = bd; f.name = FONT
        f.color.rgb = RGBColor.from_string(cl)


def title_bar(slide, txt, sub=None):
    add_text(slide, 0.5, 0.25, 12.3, 0.9, [(txt, 27, True, "1f3864")])
    if sub: add_text(slide, 0.5, 1.02, 12.3, 0.5, [(sub, 14, False, "7f7f7f")])


def pic(slide, path, l, t, w): slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))


def add_table(slide, data, l, t, w, h, col_w=None, fs=12, verdict_col=None, vmap=None):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    if col_w:
        for i, cw in enumerate(col_w): tbl.columns[i].width = Inches(cw)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c); cell.text = str(data[r][c])
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("1f3864" if r == 0 else ("eef2f8" if r % 2 else "ffffff"))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(fs); run.font.name = FONT; run.font.bold = (r == 0)
                    col = "ffffff" if r == 0 else "222222"
                    if verdict_col is not None and c == verdict_col and r > 0 and vmap:
                        col = vmap.get(data[r][c], "222222")
                        run.font.bold = True
                    run.font.color.rgb = RGBColor.from_string(col)
    return tbl


# 1 title
s = prs.slides.add_slide(BLANK)
add_text(s, 0.8, 2.0, 11.7, 2, [
    ("推理时思维链错误检测:隐藏层激活的几何信号", 33, True, "1f3864"),
    ("从激活内禀维度到「到健康推理流形的距离」", 19, False, "2f6fb0")])
add_text(s, 0.8, 4.5, 11.7, 1.8, [
    ("模型:Meta-Llama-3.1-8B-Instruct", 15, False, "444444"),
    ("数据:GSM8K 同题多采样,每题 12 条解", 15, False, "444444"),
    ("结论分三档标注:已确立 / 候选待验证 / 已证伪;数字均取自原始统计文件", 13, False, "c0392b")])

# 2 question & data
s = prs.slides.add_slide(BLANK); title_bar(s, "一、研究问题与数据")
add_text(s, 0.5, 1.5, 12.3, 0.6, [("研究问题:只看模型内部激活,能否判断一条推理链最终会不会答错", 17, True, "1f3864")])
add_table(s, [["项", "内容"],
              ["数据集", "GSM8K 测试集"],
              ["采样", "每道题采样 12 条解"],
              ["判对错", "黄金答案数值匹配,不使用大模型判别"],
              ["激活向量", "每步 token 隐状态加权聚合为一个向量,再对选定层取平均"],
              ["主评测", "题内配对 AUROC"],
              ["对照评测", "跨题 AUROC"]],
          0.5, 2.3, 9.0, 3.6, col_w=[2.2, 6.8], fs=14)
add_text(s, 0.5, 6.1, 12.3, 0.8, [
    ("同题内比较把题目难度固定住;跨题比较会把难度混进结果", 14, True, "c0392b")])

# 3 methods how-computed
s = prs.slides.add_slide(BLANK); title_bar(s, "二、方法:每个怎么算")
add_table(s, [["方法", "怎么算", "输出"],
              ["参与度 PR", "PR 等于 各分量平方和的平方 除以 各分量四次方和", "激活了几维"],
              ["激活熵 AE", "激活分量谱熵的指数", "激活分散度"],
              ["训练线性探针", "激活向量到对错标签,逻辑回归加 L2 正则,按题留出 5 折", "判别分数"],
              ["Mahalanobis 距离", "到健康均值的逐维标准化平方距离,即 各维 差的平方除以方差 再求和", "离群程度"],
              ["题内配对 AUROC", "同题内错误链分数高于正确链分数的配对占比", "检测力"],
              ["跨题 AUROC", "所有链不分题混合排序的 AUROC", "对照检测力"]],
          0.5, 1.55, 12.3, 4.4, col_w=[2.8, 7.5, 2.0], fs=13)
add_text(s, 0.5, 6.2, 12.3, 0.6, [("健康均值指正确链激活向量的均值;按题留出指训练与测试不共享同一道题", 12, False, "7f7f7f")])

# 4 evaluation settings
s = prs.slides.add_slide(BLANK); title_bar(s, "三、评测设置")
add_table(s, [["提示形式", "判分规则", "正确链", "错误链", "对照题数"],
              ["自定义零样本", "宽松", "1886", "489", "117"],
              ["标准五样本", "宽松", "2118", "528", "122"],
              ["标准五样本", "严格", "1756", "890", "204"]],
          0.5, 1.7, 10.5, 2.6, col_w=[2.8, 2.0, 1.9, 1.9, 1.9], fs=14)
add_text(s, 0.5, 4.6, 12.3, 1.6, [
    ("宽松判分:取回复最后一个数字与答案比较,相等即算正确", 14, False, "222222"),
    ("严格判分:要求回复含「#### 数字」格式行且数值精确匹配,否则算错误", 14, False, "222222"),
    ("标准五样本提示加严格判分最接近学界评测,也最干净,后续核心结果以此为准", 14, True, "c0392b")])

# 5 probe + inflation
s = prs.slides.add_slide(BLANK); title_bar(s, "结果一:训练探针揭示难度膨胀")
pic(s, figs[0], 0.4, 1.55, 6.7)
add_text(s, 7.4, 1.5, 5.6, 0.6, [("设置:自定义零样本提示,宽松判分", 13, False, "7f7f7f")])
add_table(s, [["评测方式", "AUROC"],
              ["跨题随机划分", "0.898"],
              ["按题留出整体排序", "0.715"],
              ["按题留出同题配对", "0.715"],
              ["链长基线", "0.585"]],
          7.4, 2.1, 5.4, 2.6, col_w=[3.7, 1.7], fs=13)
add_text(s, 7.4, 5.1, 5.6, 1.6, [
    ("跨题比题内高出 0.18", 15, True, "c0392b"),
    ("这 0.18 来自探针记住了题目本身,不是真错误检测", 13, False, "222222")])

# 6 emergence
s = prs.slides.add_slide(BLANK); title_bar(s, "结果二:错误检测力沿推理链后移")
pic(s, figs[1], 0.4, 1.55, 6.7)
add_text(s, 7.4, 1.5, 5.6, 0.6, [("方法:每步按位置归一分五段,各段单独训练探针", 12, False, "7f7f7f")])
add_table(s, [["链中位置", "题内配对", "跨题"],
              ["0.0–0.2", "0.554", "0.833"],
              ["0.4–0.6", "0.625", "0.844"],
              ["0.8–1.0", "0.687", "0.896"]],
          7.4, 2.1, 5.4, 1.9, col_w=[2.0, 1.8, 1.6], fs=13)
add_text(s, 7.4, 4.5, 5.6, 2.0, [
    ("题内信号从前段近随机升到后段 0.69", 14, True, "1f3864"),
    ("错误在推理后段才显现", 13, False, "222222"),
    ("跨题全程高且平,反映的是题目难度", 13, False, "222222")])

# 7 convergence
s = prs.slides.add_slide(BLANK); title_bar(s, "结果三:多种方法收敛到 0.68 至 0.76")
pic(s, figs[2], 0.4, 1.55, 6.7)
add_text(s, 7.4, 1.5, 5.6, 0.6, [("设置:自定义零样本提示,宽松判分,题内配对", 12, False, "7f7f7f")])
add_table(s, [["方法", "题内配对 AUROC"],
              ["参与度 PR", "0.60"], ["激活熵 AE", "0.58"], ["Mahalanobis", "0.65"],
              ["流形外能量 SPE", "0.68"], ["训练线性探针", "0.71"],
              ["低秩 25 维探针", "0.73"], ["多信号集成", "0.76"]],
          7.4, 2.05, 5.4, 3.4, col_w=[3.3, 2.1], fs=12)
add_text(s, 7.4, 5.7, 5.6, 1.0, [
    ("十余种方法都落在同一区间", 14, True, "1f3864"),
    ("这是这套激活表示的信息上限,不是算法不够", 12, False, "222222")])

# 8 mahalanobis 0.83
s = prs.slides.add_slide(BLANK); title_bar(s, "结果四:严格判分加标准提示下 Mahalanobis 达 0.830")
pic(s, figs[3], 0.4, 1.55, 6.7)
add_text(s, 7.4, 1.5, 5.6, 0.5, [("方法:Mahalanobis 距离,题内配对", 12, False, "7f7f7f")])
add_table(s, [["设置", "中层", "全层", "深层"],
              ["自定义零样本 宽松", "0.628", "0.646", "0.645"],
              ["标准五样本 宽松", "0.657", "0.647", "0.639"],
              ["标准五样本 严格", "0.830", "0.815", "0.811"]],
          7.4, 2.0, 5.6, 2.2, col_w=[2.6, 1.0, 1.0, 1.0], fs=12)
add_text(s, 7.4, 4.5, 5.6, 2.0, [
    ("只有提示与判分都对齐时升到 0.830", 14, True, "c0392b"),
    ("三个层带一致", 13, False, "222222"),
    ("链长干扰消失,链长题内 AUROC 为 0.486", 13, False, "222222")])

# 9 feature survival + validation
s = prs.slides.add_slide(BLANK); title_bar(s, "结果五:特征筛选与三项待验证", "设置:标准五样本提示,严格判分")
pic(s, figs[4], 0.4, 1.6, 6.6)
add_table(s, [["待验证", "现状", "做法"],
              ["是否在测格式失败", "错误链中 41% 是答对但未按格式收尾", "只取格式合规链重算"],
              ["是否过拟合", "当前用全体链估均值", "改为只用训练折正确链并按题留出"],
              ["是否只在单一模型成立", "仅 Llama 与 GSM8K", "换 Qwen 数学模型与 MATH 数据集"]],
          7.2, 1.9, 5.9, 3.2, col_w=[2.3, 2.4, 1.2], fs=11.5)
add_text(s, 7.2, 5.3, 5.9, 1.4, [
    ("24 个特征仅 3 个 Mahalanobis 存活", 14, True, "1f3864"),
    ("原假设错误激活更弥散不成立,参与度反低为 0.396", 12, False, "c0392b")])

# 10 hypotheses table
s = prs.slides.add_slide(BLANK); title_bar(s, "四、假说验证一览")
data = [["假说", "判定", "证据"],
        ["错误推理激活更弥散", "已证伪", "严格判分下参与度反低,0.396"],
        ["跨题评测被难度抬高", "已确立", "0.898 比 0.715 高 0.18"],
        ["错误检测力沿链后移", "已确立", "题内从 0.554 升到 0.687"],
        ["难度与错误是不同方向", "已确立", "夹角余弦 0.24,投影掉难度后不降"],
        ["错误方向对应犹豫回溯语义", "已确立", "词表投影加洗牌对照"],
        ["错误链偏离健康激活区", "候选", "Mahalanobis 0.830,三项待验证"],
        ["错误与正确占据不同子空间", "未验证", "需主夹角与集合外能量实验"]]
vmap = {"已确立": "2e8b57", "已证伪": "c0392b", "候选": "e08e0b", "未验证": "7f7f7f"}
add_table(s, data, 0.5, 1.6, 12.3, 4.8, col_w=[4.2, 2.0, 6.1], fs=13, verdict_col=1, vmap=vmap)

# 11 next steps
s = prs.slides.add_slide(BLANK); title_bar(s, "五、下一步:从检测到干预")
add_text(s, 0.5, 1.6, 4.1, 5.2, [
    ("先核验", 17, True, "1f3864"),
    ("剔除格式失败后重算", 13, False, "222222"),
    ("改用按题留出且只用正确链估均值", 13, False, "222222"),
    ("对角距离与全协方差距离对比", 13, False, "222222")])
add_text(s, 4.7, 1.6, 4.1, 5.2, [
    ("做成检测器", 17, True, "1f3864"),
    ("给出判定阈值与风险覆盖曲线", 13, False, "222222"),
    ("与困惑度 自洽性 模型自评比较", 13, False, "222222"),
    ("换 Qwen 数学模型与 MATH 数据集", 13, True, "c0392b")])
add_text(s, 8.9, 1.6, 4.0, 5.2, [
    ("用于干预", 17, True, "1f3864"),
    ("诊断型,推荐", 14, True, "2e8b57"),
    ("低置信度时弃答或转交更强模型", 12, False, "222222"),
    ("用作验证器对多解重排", 12, False, "222222"),
    ("因果型,谨慎", 14, True, "c0392b"),
    ("直接编辑激活已被现有工作证明无效", 12, False, "222222"),
    ("仅列为远期探索", 12, False, "222222")])

path = os.path.join(OUT, "项目进展_PPT_2026-06-08.pptx"); prs.save(path)
print("PPT saved ->", path, "slides:", len(prs.slides._sldIdLst))
