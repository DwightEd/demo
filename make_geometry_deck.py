"""Build the geometry deck (honest version): two named directional metrics, plain-language
construction, why it works, theory, results WITH difficulty control, competitors as context.
Run where python-pptx is installed:  pip install python-pptx; python make_geometry_deck.py
Output: 项目进展_PPT_几何指标_2026-06-24.pptx . Numbers are the difficulty/length-controlled honest ones."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

DST = "项目进展_PPT_几何指标_2026-06-24.pptx"

# kinds: ("section", title) | ("content", title, [lines], table_or_None);  table = (header, rows)
S = [
    ("section", "推理步级错误检测：两个方向几何指标\n方向集中度 κ 与方向广度 ρ｜模型 Llama-3.1-8B｜数据 ProcessBench"),

    ("content", "我们做的事", [
        "给一条多步数学推理，判断哪一步是第一个出错的步。",
        "只看模型内部的隐藏层激活，不看输出文字。",
        "方法就是两个可解释、不需要训练的方向几何指标：方向集中度 κ 和方向广度 ρ。"], None),

    ("content", "指标一：方向集中度 κ（怎么算）", [
        "一个推理步里有若干个 token，每个 token 在第 14 层有一个 4096 维向量。",
        "每个向量除以自己的长度，变成只看方向、不看长短的单位向量。",
        "把这些单位向量加权平均（靠后的 token 权重大一点），得到一个平均方向。",
        "量这个平均方向的长度就是 κ，范围 0 到 1。",
        "都指一个方向，平均向量长，κ 接近 1；指得乱、互相抵消，κ 接近 0。",
        "含义：κ 高是这一步方向一致、推理聚焦；κ 低是方向发散、可能在乱走。"], None),

    ("content", "指标二：方向广度 ρ（怎么算）", [
        "还是这一步的那些单位向量。",
        "算它们的散布矩阵，每个向量和自己的外积求平均，描述这堆方向铺在哪些方向上。",
        "求散布矩阵的特征值，归一化到加起来等于 1。",
        "用特征值算一个有效方向数 ρ：能量全在一个方向，ρ 等于 1；平摊到很多方向，ρ 很大。",
        "含义：ρ 小是方向集中在少数几个，结构化；ρ 大是铺得到处都是，各向同性发散。"], None),

    ("content", "两个指标的分工", [
        "κ 只看有没有一个主方向，是一阶矩。",
        "ρ 看一共有几个方向，是二阶矩。",
        "κ 把乱发散和沿几个方向铺开混在一起，两种都是低 κ；ρ 能把它们分开。"], None),

    ("content", "为什么有用：错误步方向更发散，κ 更低", [
        "把步按 κ 从低到高分十段，看每段的错误率，单调下降，没有反弹。",
        "说明 κ 是一个可以挖掘的方向，不是噪声。"],
     (["κ 从低到高", "最低段", "2", "3", "4", "5", "6", "7", "8", "9", "最高段"],
      [["gsm8k 错误率", "0.22", "0.25", "0.21", "0.13", "0.09", "0.04", "0.04", "0.06", "0.08", "0.04"],
       ["omnimath 错误率", "0.21", "0.16", "0.16", "0.14", "0.15", "0.10", "0.14", "0.11", "0.11", "0.09"]])),

    ("content", "理论依据：方向统计", [
        "κ 是 von Mises-Fisher 分布集中度的充分统计量，方向集中这一维已经被它抓全。",
        "这解释了为什么各种动态、相对、形状的变体在 κ 之上都加不出东西。",
        "ρ 对应更弱假设的 Bingham 分布，用二阶矩，能看到 κ 看不到的发散结构。",
        "诚实说明：这套分布是模型，它的核心预言被实验间接支持，但还没做直接的分布拟合检验。"], None),

    ("content", "结果一：不控难度时，几何和熵差不多", [
        "单个 κ 对单个熵指标 EDIS，混在一起排序的 AUROC。",
        "两者基本打平，几何并没有单点碾压熵。",
        "重要：这是不控难度的口径，下一页说明问题在哪。"],
     (["子集", "单个 κ", "单个 EDIS"],
      [["gsm8k", "0.772", "0.719"], ["math", "0.703", "0.717"],
       ["omnimath", "0.702", "0.754"], ["olympiad", "0.703", "0.753"]])),

    ("content", "难度控制怎么做", [
        "ProcessBench 每道题只有一条解，正确链和错误链来自不同的题。",
        "不控难度时，模型只要认出这道题难就能蒙对，分数被难度抬高。",
        "难度控制就是在同一道题内部比：拿首错步和它自己前面的正确步比，难度固定。",
        "长度控制就是比之前先去掉步有多长的影响，因为错误步往往更长。"], None),

    ("content", "结果二：控住难度和长度后，几何很弱，熵略强", [
        "在同一道题内、去掉长度影响后，定位首错步的能力。",
        "κ 只剩 0.55 左右，勉强超过随机的 0.5。",
        "熵在难任务上反而比几何更高。",
        "二阶矩 ρ 在难任务上有一个干净的小增量，加了非线性长度项后仍然成立。"],
     (["子集", "κ 控全后", "熵 控全后"],
      [["gsm8k", "0.574", "0.537"], ["math", "0.560", "0.583"],
       ["omnimath", "0.542", "0.609"], ["olympiad", "0.552", "0.620"]])),

    ("content", "竞品只能当背景，不能直接比", [
        "其他工作报的都是不控难度的池化数，而且模型、数据、协议各不相同。",
        "GeoReason 在 ProcessBench 上报 0.91，但那是同分布内最乐观的数，跨模型掉到 0.75，也没说用哪个模型。",
        "Streaming 报 0.88，但在多跳问答数据上、用知道答案的裁判标的步、也没控难度。",
        "Hidden Error 报 0.95，自己承认在同题内掉到 0.70。",
        "结论：这些数不能和我们直接比，要比必须在同一模型同一协议下重跑，我们没做。"], None),

    ("content", "诚实的结论与我们的贡献", [
        "不控难度时，几何和熵都能到 0.70 到 0.80，和这个领域一个量级。",
        "但这些分数大半来自难度和长度，控住之后真实信号只剩 0.55，熵在难任务还更强。",
        "二阶矩 ρ 是一个干净的小增量，扛住了非线性长度检验。",
        "我们真正的贡献是难度加长度的双控协议，揭穿这类内部信号的虚高；以及两个可解释、不用训练的方向几何指标。"], None),
]


def add(prs, kind, *a):
    if kind == "section":
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(8.8), Inches(2.2)).text_frame
        tf.word_wrap = True
        for i, line in enumerate(a[0].split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.size = Pt(26 if i == 0 else 15); p.runs[0].font.bold = (i == 0)
        return
    title, lines, table = a
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9)).text_frame
    t.word_wrap = True; t.text = title
    t.paragraphs[0].runs[0].font.size = Pt(24); t.paragraphs[0].runs[0].font.bold = True
    y = 1.3
    for b in lines:
        p = s.shapes.add_textbox(Inches(0.6), Inches(y), Inches(8.8), Inches(0.55)).text_frame
        p.word_wrap = True; p.text = b; p.paragraphs[0].runs[0].font.size = Pt(15); y += 0.55
    if table:
        hdr, rows = table; nr, nc = len(rows) + 1, len(hdr)
        gt = s.shapes.add_table(nr, nc, Inches(0.6), Inches(y + 0.15), Inches(8.8), Inches(0.36 * nr)).table
        for j, h in enumerate(hdr):
            gt.cell(0, j).text = h
        for i, r in enumerate(rows):
            for j, v in enumerate(r):
                gt.cell(i + 1, j).text = str(v)
        for rr in range(nr):
            for cc in range(nc):
                pr = gt.cell(rr, cc).text_frame.paragraphs[0]
                if pr.runs:
                    pr.runs[0].font.size = Pt(12)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    for item in S:
        add(prs, item[0], *item[1:])
    prs.save(DST)
    print(f"saved {DST} ({len(S)} slides)")


if __name__ == "__main__":
    main()
